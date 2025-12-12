import xlwings as wx
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

try:
    from .GetData import Futam, Horses
except:
    from GetData import Futam, Horses


class MakePPT:
    def __init__(self,drivers,titles):
        self.rome_num = ["I","II","III","IV","V","VI","VII","VIII","IX","X","XI","XII","XIII","XIV"]
        self.titles = titles
        self.drivers = drivers

        if not os.path.isdir("ppt"):
            os.makedirs("ppt")

        for title in self.titles:
            ppt = Presentation()
            
            ppt.slide_width = Inches(10)  # Set width
            ppt.slide_height = Inches(7.5)  # Set height
            slide_layout = ppt.slide_layouts[6]

            self.slide1(ppt,slide_layout,title)
            self.slide2(ppt,slide_layout,title)

            if(title.id==0):
                self.slide3(ppt,slide_layout,title)
                self.slide4(ppt,slide_layout,title,True)

            elif(self.titles[-1].id == title.id):
                self.slide3(ppt,slide_layout,self.titles[int(title.id)-1])
                self.slide4(ppt,slide_layout,self.titles[int(title.id)-1])
                self.slide3(ppt,slide_layout,title,True)
                self.slide4(ppt,slide_layout,title,True)

            else:
                self.slide3(ppt,slide_layout,self.titles[int(title.id)-1])
                self.slide4(ppt,slide_layout,self.titles[int(title.id)-1])

            self.slide5(ppt,slide_layout)

            file_name = ""
            file_name =f".\\ppt\\{title.daily}. futam.pptx"

            
            ppt.save(file_name)

        self.run_vba_macro()
        
    def slide1(self,ppt,slide_layout,futam):
        slide1 = ppt.slides.add_slide(slide_layout)
        #self.set_slide_duration(slide1, 5) 
        #slide1.slide_show.transition.duration = 50
        slide1bc = slide1.background.fill
        slide1bc.solid()
        slide1bc.fore_color.rgb = RGBColor(0, 0, 0)
        
        title_box = slide1.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1.4))
        title_frame = title_box.text_frame
        
        title = title_frame.add_paragraph()
        #print(futam)
        title_str = f"{futam.daily}. {futam.title}".strip()

        title.text = title_str
        if   (len(title_str) < 39 ): title.font.size = Pt(45)
        elif (len(title_str) <= 58): title.font.size = Pt(40)
        else:                        title.font.size = Pt(36)
        title.font.bold = True
        title.font.color.rgb = RGBColor(255, 229, 121)
        
        title.alignment = PP_ALIGN.CENTER
        title_frame.word_wrap = True 
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        text_box = slide1.shapes.add_textbox(Inches(4.93), Inches(1.5), Inches(5.07), Inches(0.7))
        text_frame = text_box.text_frame
        
        text = text_frame.add_paragraph()
        text.text = f"Pálya: 11 Kincsem Park"
        text.font.size = Pt(36)
        text.font.bold = True
        text.font.color.rgb = RGBColor(255, 229, 121)
        
        text.alignment = PP_ALIGN.RIGHT
        text_frame.word_wrap = True 
        text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        
        data_box = slide1.shapes.add_textbox(Inches(0), Inches(2.53), Inches(3.69), Inches(1.31))
        data_frame = data_box.text_frame
        
        data_text = data_frame.add_paragraph()
        data_text.text = f"{futam.dist} m\n{futam.start}"

        data_text.font.size = Pt(36)
        data_text.font.bold = True
        data_text.font.color.rgb = RGBColor(255, 229, 121)
        
        data_text.alignment = PP_ALIGN.LEFT
        data_frame.word_wrap = True 
        data_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        
        time_box = slide1.shapes.add_textbox(Inches(5.52), Inches(2.53), Inches(1.51), Inches(1.31))
        time_frame = time_box.text_frame
        
        time_text = time_frame.add_paragraph()
        time_text.text = f"Start:\n{futam.time}"

        time_text.font.size = Pt(36)
        time_text.font.bold = True
        time_text.font.color.rgb = RGBColor(255, 255, 255)
        
        time_text.alignment = PP_ALIGN.LEFT
        time_frame.word_wrap = True 
        time_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        
        
        slide1.shapes.add_picture("clock.jpeg", Inches(7.47), Inches(2.63), Inches(1.17), Inches(1.13))
        
        text1_box = slide1.shapes.add_textbox(Inches(0), Inches(4.94), Inches(3.3), Inches(0.71))
        text1_frame = text1_box.text_frame
        
        text1 = text1_frame.add_paragraph()
        text1.text = f"Véleményünk:"
        text1.font.size = Pt(36)
        text1.font.bold = True
        text1.font.color.rgb = RGBColor(255, 255, 255)
        
        text1.alignment = PP_ALIGN.LEFT
        text1_frame.word_wrap = True 
        text1_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        
        opinion_box = slide1.shapes.add_textbox(Inches(0), Inches(5.34), Inches(10), Inches(1.92))
        opinion_frame = opinion_box.text_frame
        
        opinion_text = opinion_frame.add_paragraph()
        opinion_text.text = f"{futam.opinion}"

        opinion_text.font.size = Pt(30)
        opinion_text.font.bold = False
        opinion_text.font.color.rgb = RGBColor(255, 255, 255)
        
        time_text.alignment = PP_ALIGN.LEFT
        opinion_frame.word_wrap = True 
        opinion_frame.vertical_anchor = MSO_ANCHOR.TOP

    def slide2(self,ppt,slide_layout,futam):
        drivers_list = [driver for driver in self.drivers if driver.Fnum == futam.id]

        slide2 = ppt.slides.add_slide(slide_layout)
        #slide2.slide_show.transition.duration = 50
        slide2bc = slide2.background.fill
        slide2bc.solid()
        slide2bc.fore_color.rgb = RGBColor(0, 0, 0)
        
        horse_box = slide2.shapes.add_textbox(Inches(0), Inches(0), Inches(5), Inches(7.5))
        horse_frame = horse_box.text_frame
        
        driver_box = slide2.shapes.add_textbox(Inches(5), Inches(0), Inches(5), Inches(7.5))
        driver_frame = driver_box.text_frame
        
        for row in drivers_list:
            
            #horse side

            horse = horse_frame.add_paragraph()
            horse.text = f"{row.Hnum}. {row.Hname.upper()}"
            #print(f"{row.horse_number}. {row.horse_name.upper()}")
            horse.font.size = Pt(32)
            horse.font.bold = True
            horse.font.color.rgb = RGBColor(255, 229, 121)
            
            horse.alignment = PP_ALIGN.LEFT
            horse_frame.word_wrap = True 
            horse_frame.vertical_anchor = MSO_ANCHOR.TOP
            
            #driver side
            
            driver = driver_frame.add_paragraph()
            driver.text = f"{row.DJname}"
            driver.font.size = Pt(32)
            driver.font.bold = False
            driver.font.color.rgb = RGBColor(255, 255, 255)
            
            driver.alignment = PP_ALIGN.LEFT
            driver_frame.word_wrap = True 
            driver_frame.vertical_anchor = MSO_ANCHOR.TOP

    def slide3(self,ppt,slide_layout,futam,hide=False):
        slide3 = ppt.slides.add_slide(slide_layout)
        #slide3.slide_show.transition.duration = 50
        slide3bc = slide3.background.fill
        slide3bc.solid()
        slide3bc.fore_color.rgb = RGBColor(0, 0, 0)
        

        
        title_box = slide3.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1.4))
        title_frame = title_box.text_frame
        
        title = title_frame.add_paragraph()
        
        title_str = f"{futam.daily}. {futam.title}".strip()

        title.text = title_str
        if   (len(title_str) < 39 ): title.font.size = Pt(45)
        elif (len(title_str) <= 58): title.font.size = Pt(40)
        else:                        title.font.size = Pt(36)
        title.font.bold = True
        title.font.color.rgb = RGBColor(255, 229, 121)
        
        title.alignment = PP_ALIGN.CENTER
        title_frame.word_wrap = True 
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            

            
        text_box = slide3.shapes.add_textbox(Inches(0), Inches(1.4), Inches(6.74), Inches(1.72))
        text_frame = text_box.text_frame
        
        text = text_frame.add_paragraph()
        text.text = f"Pálya: 10 Kincsem Park\nBefutási sorrend:"
        text.font.size = Pt(48)
        text.font.bold = True
        text.font.color.rgb = RGBColor(255, 229, 121)
        
        text.alignment = PP_ALIGN.LEFT
        text_frame.word_wrap = True 
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        order_box = slide3.shapes.add_textbox(Inches(0), Inches(3.37), Inches(10), Inches(2.52))
        order_frame = order_box.text_frame
        
        order = order_frame.add_paragraph()
        order.text = f"I.\nII.\nIII."
        order.font.size = Pt(48)
        order.font.bold = True
        order.font.color.rgb = RGBColor(255, 255, 255)
        
        order.alignment = PP_ALIGN.LEFT
        order_frame.word_wrap = True 
        order_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        time_box = slide3.shapes.add_textbox(Inches(0), Inches(5.9), Inches(3.39), Inches(0.91))
        time_frame = time_box.text_frame
        
        time_text = time_frame.add_paragraph()
        time_text.text = f"Idő:"

        time_text.font.size = Pt(48)
        time_text.font.bold = True
        time_text.font.color.rgb = RGBColor(255, 229, 121)
        
        time_text.alignment = PP_ALIGN.LEFT
        time_frame.word_wrap = True 
        time_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        
    def slide4(self,ppt,slide_layout,futam,hide=False):
        slide4 = ppt.slides.add_slide(slide_layout)
        #slide4.slide_show.transition.duration = 50
        slide4bc = slide4.background.fill
        slide4bc.solid()
        slide4bc.fore_color.rgb = RGBColor(0, 0, 0)
        
        title_box = slide4.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1.4))
        title_frame = title_box.text_frame
        title = title_frame.add_paragraph()

            
        title_str = f"{futam.daily}. {futam.title}".strip()

        title.text = title_str
        if   (len(title_str) < 39 ): title.font.size = Pt(45)
        elif (len(title_str) <= 58): title.font.size = Pt(40)
        else:                        title.font.size = Pt(36)
        title.font.bold = True
        title.font.color.rgb = RGBColor(255, 229, 121)
        
        title.alignment = PP_ALIGN.CENTER
        title_frame.word_wrap = True 
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
        
        text_box = slide4.shapes.add_textbox(Inches(0), Inches(1.4), Inches(10), Inches(1.72))
        text_frame = text_box.text_frame
        
        text = text_frame.add_paragraph()
        text.text = f"Pálya: 10 Kincsem Park\nBefutási sorrend:"
        text.font.size = Pt(48)
        text.font.bold = True
        text.font.color.rgb = RGBColor(255, 229, 121)
        
        text.alignment = PP_ALIGN.LEFT
        text_frame.word_wrap = True 
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        dividend_box = slide4.shapes.add_textbox(Inches(0), Inches(3.37), Inches(10), Inches(2.52))
        dividend_frame = dividend_box.text_frame
        
        dividend = dividend_frame.add_paragraph()
        dividend.text = f"Tét:					1\nHely:					1\nBefutó:    		1\nHbefutó: 		1"
        dividend.font.size = Pt(48)
        dividend.font.bold = True
        dividend.font.color.rgb = RGBColor(255, 255, 255)
        
        dividend.alignment = PP_ALIGN.LEFT
        dividend_frame.word_wrap = True 
        dividend_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
    def slide5(self,ppt,slide_layout):
        slide5 = ppt.slides.add_slide(slide_layout)
        #slide5.slide_show.transition.duration = 50
        slide5bc = slide5.background.fill
        slide5bc.solid()
        slide5bc.fore_color.rgb = RGBColor(0, 0, 0)
        
        title_box = slide5.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(2.3))
        title_frame = title_box.text_frame
        
        title = title_frame.add_paragraph()
        title.text = f"Nem hivatalos\n  befutási sorrend:"
        title.font.size = Pt(88)
        title.font.bold = True
        title.font.color.rgb = RGBColor(255, 229, 121)
        
        title.alignment = PP_ALIGN.LEFT
        title_frame.word_wrap = True 
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        text_box = slide5.shapes.add_textbox(Inches(0), Inches(3.75), Inches(10), Inches(2.21))
        text_frame = text_box.text_frame
        
        text = text_frame.add_paragraph()
        text.text = f" –  – "
        text.font.size = Pt(125)
        text.font.bold = True
        text.font.color.rgb = RGBColor(255, 255, 255)
        
        text.alignment = PP_ALIGN.CENTER
        text_frame.word_wrap = True 
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    

    def run_vba_macro(self):
        wb = wx.Book("add macro.xlsm")
        macro1 = wb.macro("Module1.SetPPTSlidesFromFolderAndExitExcel")
        macro1()

if __name__ == "__main__":
    titles = []
    drivers = []
    with open("./csv/titles_data.csv",'r',encoding="utf-8") as f:
        fs = f.readline()
        for ln in f: 
            #print(ln.strip())
            titles.append(Futam(ln))

    with open("./csv/drivers_data.csv",'r',encoding="utf-8") as f:
        fs = f.readline()
        for ln in f: drivers.append(Horses(ln))

    MakePPT(drivers,titles)
    print("Id;Daily;Title;Distance;Start time;Start type;Opinion")
    for i in titles: print(i)

